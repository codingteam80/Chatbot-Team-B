from pptx import Presentation

from ingestion.document import ParsedDocument
from ingestion.loaders.base_loader import BaseLoader


class PPTXLoader(BaseLoader):

    def load(self, file_path):

        prs = Presentation(file_path)

        documents = []

        for slide_number, slide in enumerate(
            prs.slides,
            start=1
        ):

            texts = []

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    text = shape.text.strip()

                    if text:

                        texts.append(text)

            if not texts:
                continue

            documents.append(

                ParsedDocument(
                    text="\n".join(texts),
                    slide_number=slide_number
                )

            )

        print(
            f"[PPTX] Slides created: {len(documents)}"
        )

        return documents